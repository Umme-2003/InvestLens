# backend/scripts/extract_text.py
import sys
import io # Import the io module
from pptx import Presentation
import pytesseract
from PIL import Image

def extract_text_from_pptx(file_stream): # It now accepts a stream
    # The pptx library can open a file-like object directly
    prs = Presentation(file_stream)
    full_text = []

    for slide in prs.slides:
        slide_text = []
        for shape in slide.shapes:
            if hasattr(shape, "text"):
                slide_text.append(shape.text)
        
        for shape in slide.shapes:
            if shape.shape_type == 13:
                try:
                    image = shape.image
                    image_bytes = image.blob
                    img = Image.open(io.BytesIO(image_bytes))
                    ocr_text = pytesseract.image_to_string(img)
                    slide_text.append(ocr_text)
                except Exception as e:
                    pass

        full_text.append('\n'.join(slide_text))

    return "\n--- SLIDE BREAK ---\n".join(full_text)

if __name__ == "__main__":
    sys.stdout.reconfigure(encoding='utf-8')
    
    # Read the binary data from standard input
    pptx_data = sys.stdin.buffer.read()
    
    # Create an in-memory binary stream
    file_like_object = io.BytesIO(pptx_data)
    
    extracted_text = extract_text_from_pptx(file_like_object)
    print(extracted_text)